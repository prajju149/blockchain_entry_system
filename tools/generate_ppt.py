from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT = 'presentation.pptx'

prs = Presentation()

# Title slide
title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = "Blockchain-Based Entry/Exit System"
subtitle = slide.placeholders[1]
subtitle.text = f"Class Presentation\nAuthor: Your Name\nDate: {date.today().isoformat()}"

# Helper to add bullet slides
def add_bullet_slide(title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)

# Slides
add_bullet_slide("Objective", [
    "Build a secure entry/exit tracking prototype for multi-tenant housing",
    "Use QR codes + optional biometric fingerprints",
    "Record events to a tamper-evident file-backed blockchain"
])

add_bullet_slide("Key Features", [
    "QR generation and live browser scanning",
    "Per-resident profile storage (JSON)",
    "Blockchain ledger with SHA256 + HMAC-SHA256",
    "Simple Flask REST API + single file dashboard"
])

add_bullet_slide("System Architecture", [
    "Frontend: single `dashboard.html` (QR scanner, register, ledger viewer)",
    "Backend: Flask app serving API endpoints and QR creation",
    "Storage: JSON files for residents, chain, and QR assets",
    "Security: HMAC signing for ledger integrity"
])

add_bullet_slide("Demo Steps", [
    "1) Run the server locally: `python wsgi.py`",
    "2) Open dashboard: `http://127.0.0.1:5000/`",
    "3) Register a resident and scan QR to log entries",
    "4) View ledger tab to see signed events"
])

add_bullet_slide("Deployment", [
    "GitHub Actions → Heroku pipeline included",
    "Procfile + runtime.txt + gunicorn used for production",
    "See GITHUB_DEPLOYMENT.md and DEPLOYMENT.md for details"
])

add_bullet_slide("Data Flow & Blockchain", [
    "Event generated on scan → stored in ledger JSON",
    "Each block includes previous hash and HMAC signature",
    "Tamper evidence provided by chained hashes and HMAC"
])

add_bullet_slide("Next Steps", [
    "Replace biometric stubs with ML models (optional)",
    "Add authentication & HTTPS for production",
    "Migrate storage to a small DB for concurrency"
])

add_bullet_slide("Thank You / Questions", [
    "Contact: your-email@example.com",
    "Repo: include link to GitHub for project code"
])

# Save
prs.save(OUTPUT)
print(f"Presentation created: {OUTPUT}")
